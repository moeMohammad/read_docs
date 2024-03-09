import fitz
import re
import docx2txt
from pptx import Presentation


def main():
    found_kw = {}
    kw = get_kw("keywords")
    text = extract_text("sample.pdf")
    for word in kw:
        if is_word_in_text(word, str(text)):
            found_kw[word] = entropy_kw(word, str(text))
    print(found_kw)


def get_kw(file_name):
    kw = []
    with open(file_name) as file:
        for line in file:
            kw.append(line.lower().strip())
    return kw


def extract_text(file):
    if file.endswith(".pdf"):
        doc = fitz.open(file)
        text = ""
        for page in doc:
            text += page.get_text()
        text = text.splitlines()
    elif file.endswith(".docx"):
        text = docx2txt.process(file)
    elif file.endswith(".pptx"):
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        text = "\n".join(text_runs)
    
    return text


def is_word_in_text(word : str, text : str) -> bool:
    """
    Check if a word is in a text.

    Parameters
    ----------
    word : str
    text : str

    Returns
    -------
    bool : True if word is in text, otherwise False.

    Examples
    --------
    is_word_in_text("Python", "python is awesome.")
    True

    is_word_in_text("Python", "camelCase is pythonic.")
    False

    is_word_in_text("Python", "At the end is Python")
    True
    """
    if "+" in word:
        word = re.escape(word)
    pattern = r'(^|[^\w]){}([^\w]|$)'.format(word)
    pattern = re.compile(pattern, re.IGNORECASE)
    matches = re.search(pattern, text)
    return bool(matches)

def entropy_kw(word, string):
    if "+" in word:
        word = re.escape(word)
    pattern = r'(^|[^\w]){}([^\w]|$)'.format(word)
    matches = re.findall(pattern, string, re.IGNORECASE)
    entropy = 1/len(matches)
    if entropy.is_integer():
        return int(entropy)
    else:
        return round(entropy, 3)
    
if __name__ == "__main__":
    main()
